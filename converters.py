"""
Convert file bytes (pdf/docx/pptx) to Markdown, entirely in memory.
"""

import io


def convert_to_markdown(data: bytes, ext: str, filename: str) -> str:
    ext = ext.lower()
    if ext == "pdf":
        return _pdf_to_md(data)
    if ext == "docx":
        return _docx_to_md(data)
    if ext == "pptx":
        return _pptx_to_md(data)
    raise ValueError(f"Unsupported extension: {ext} ({filename})")


# ── PDF ──────────────────────────────────────────────────────

def _pdf_to_md(data: bytes) -> str:
    import fitz  # PyMuPDF
    import pymupdf4llm

    doc = fitz.open(stream=data, filetype="pdf")
    try:
        return pymupdf4llm.to_markdown(doc)
    finally:
        doc.close()


# ── DOCX ─────────────────────────────────────────────────────

def _docx_to_md(data: bytes) -> str:
    import docx

    document = docx.Document(io.BytesIO(data))
    lines = []

    for block in _iter_docx_blocks(document):
        if block["type"] == "heading":
            lines.append(f"{'#' * block['level']} {block['text']}")
        elif block["type"] == "bullet":
            lines.append(f"- {block['text']}")
        elif block["type"] == "number":
            lines.append(f"1. {block['text']}")
        elif block["type"] == "table":
            lines.append(block["text"])
        else:
            lines.append(block["text"])
        lines.append("")

    return "\n".join(lines).strip() + "\n"


def _iter_docx_blocks(document):
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    body = document.element.body
    for child in body.iterchildren():
        if child.tag == qn("w:p"):
            para = Paragraph(child, document)
            text = para.text.strip()
            if not text:
                continue
            style = (para.style.name or "").lower()
            if style.startswith("heading"):
                digits = "".join(c for c in style if c.isdigit()) or "1"
                yield {"type": "heading", "level": min(int(digits), 6), "text": text}
            elif style.startswith("list bullet"):
                yield {"type": "bullet", "text": text}
            elif style.startswith("list number"):
                yield {"type": "number", "text": text}
            else:
                yield {"type": "paragraph", "text": text}
        elif child.tag == qn("w:tbl"):
            table = Table(child, document)
            rows = [[cell.text for cell in row.cells] for row in table.rows]
            yield {"type": "table", "text": _rows_to_md_table(rows)}


# ── PPTX ─────────────────────────────────────────────────────

def _pptx_to_md(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    lines = []

    for i, slide in enumerate(prs.slides, start=1):
        title_shape = slide.shapes.title
        title_id = title_shape.shape_id if title_shape else None
        title = title_shape.text_frame.text.strip() if title_shape and title_shape.has_text_frame else ""

        lines.append(f"# Слайд {i}: {title}".rstrip())
        lines.append("")

        for shape in slide.shapes:
            if shape.shape_id == title_id:
                continue
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if not text:
                        continue
                    indent = "  " * (para.level or 0)
                    lines.append(f"{indent}- {text}")
            if shape.has_table:
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                lines.append(_rows_to_md_table(rows))

        notes = _slide_notes(slide)
        if notes:
            lines.append("")
            lines.append(f"> **Нотатки доповідача:** {notes}")

        lines.append("")

    return "\n".join(lines).strip() + "\n"


def _slide_notes(slide):
    if not slide.has_notes_slide:
        return None
    text = slide.notes_slide.notes_text_frame.text.strip()
    return text or None


# ── SHARED ───────────────────────────────────────────────────

def _rows_to_md_table(rows) -> str:
    rows = [[(c or "").strip().replace("\n", " ") for c in row] for row in rows]
    if not rows:
        return ""
    header, *body = rows
    out = ["| " + " | ".join(header) + " |", "| " + " | ".join(["---"] * len(header)) + " |"]
    out += ["| " + " | ".join(r) + " |" for r in body]
    return "\n".join(out)
